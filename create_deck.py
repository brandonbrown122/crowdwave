"""
Crowdwave Accuracy Framework - Professional PowerPoint Deck
McKinsey-style presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Colors
NAVY = RgbColor(0, 51, 102)  # #003366
WHITE = RgbColor(255, 255, 255)
LIGHT_GRAY = RgbColor(248, 249, 250)
DARK_GRAY = RgbColor(102, 102, 102)
BLACK = RgbColor(34, 34, 34)
GREEN = RgbColor(40, 167, 69)
YELLOW = RgbColor(255, 193, 7)
RED = RgbColor(220, 53, 69)
LIGHT_GREEN = RgbColor(212, 237, 218)
LIGHT_YELLOW = RgbColor(255, 243, 205)
LIGHT_RED = RgbColor(248, 215, 218)

def set_font(run, size=12, bold=False, color=BLACK, name='Calibri'):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Navy background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(8.5), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_font(p.runs[0], size=40, bold=True, color=WHITE)
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(8.5), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        set_font(p.runs[0], size=18, color=WHITE)
        p.alignment = PP_ALIGN.LEFT
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.75), Inches(5), Inches(8.5), Inches(0.3))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "February 2026"
    set_font(p.runs[0], size=14, color=WHITE)
    
    return slide

def add_content_slide(prs, title, source=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title with underline
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    set_font(p.runs[0], size=18, bold=True, color=NAVY)
    
    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(9), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.25))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Crowdwave | Confidential"
    set_font(p.runs[0], size=8, color=DARK_GRAY)
    
    # Source (if provided)
    if source:
        source_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.7), Inches(9), Pt(1))
        source_line.fill.solid()
        source_line.fill.fore_color.rgb = RgbColor(224, 224, 224)
        source_line.line.fill.background()
        
        source_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(9), Inches(0.3))
        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Source: {source}"
        set_font(p.runs[0], size=8, color=DARK_GRAY)
    
    return slide

def add_metric_box(slide, left, top, width, height, value, label, color=NAVY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = RgbColor(224, 224, 224)
    
    # Value
    val_box = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.5))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(value)
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=36, color=color)
    
    # Label
    lbl_box = slide.shapes.add_textbox(left, top + Inches(0.55), width, Inches(0.4))
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=9, color=DARK_GRAY)

def add_zone_box(slide, left, top, width, height, zone_title, error, items, action, zone_color, light_color):
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = light_color
    box.line.color.rgb = zone_color
    
    # Left border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    border.fill.solid()
    border.fill.fore_color.rgb = zone_color
    border.line.fill.background()
    
    # Title
    if zone_title:
        title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.2), Inches(0.25))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = zone_title
        if p.runs:
            set_font(p.runs[0], size=11, bold=True, color=zone_color)
    
    # Error
    if error:
        error_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.35), width - Inches(0.2), Inches(0.3))
        tf = error_box.text_frame
        p = tf.paragraphs[0]
        p.text = error
        p.alignment = PP_ALIGN.CENTER
        set_font(p.runs[0], size=22, bold=True, color=NAVY)
    
    # Items
    items_top = top + Inches(0.7) if error else top + Inches(0.35)
    items_box = slide.shapes.add_textbox(left + Inches(0.15), items_top, width - Inches(0.2), Inches(0.8))
    tf = items_box.text_frame
    tf.word_wrap = True
    filtered_items = [item for item in items if item]
    for i, item in enumerate(filtered_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        if p.runs:
            set_font(p.runs[0], size=9, color=BLACK)
    
    # Action
    if action:
        action_box = slide.shapes.add_textbox(left + Inches(0.15), top + height - Inches(0.35), width - Inches(0.2), Inches(0.25))
        tf = action_box.text_frame
        p = tf.paragraphs[0]
        p.text = action
        if p.runs:
            set_font(p.runs[0], size=9, bold=True, color=zone_color)

def add_table(slide, left, top, width, data, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    
    table = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.3 * rows)).table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            
            # Header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        set_font(run, size=10, bold=True, color=WHITE)
            else:
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        set_font(run, size=10, color=BLACK)
            
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_highlight_box(slide, left, top, width, height, text, is_dark=False):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    if is_dark:
        box.fill.fore_color.rgb = NAVY
        text_color = WHITE
    else:
        box.fill.fore_color.rgb = RgbColor(232, 244, 252)
        text_color = BLACK
    box.line.fill.background()
    
    # Left accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY if not is_dark else WHITE
    accent.line.fill.background()
    
    text_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.2), height - Inches(0.1))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    set_font(p.runs[0], size=11, color=text_color)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, "Crowdwave Accuracy Framework", "Calibrated AI predictions with documented, predictable accuracy")

# Slide 2: Key Metrics
slide = add_content_slide(prs, "Calibration reduces prediction error by 79%, enabling reliable AI-simulated research", 
                          "Crowdwave validation (27 test cases); ForecastBench 2025")

# Metrics row
add_metric_box(slide, Inches(0.7), Inches(1.2), Inches(2), Inches(1), "79%", "Error reduction\nvs. naive LLM")
add_metric_box(slide, Inches(2.9), Inches(1.2), Inches(2), Inches(1), "1.9", "Mean absolute\nerror (points)")
add_metric_box(slide, Inches(5.1), Inches(1.2), Inches(2), Inches(1), "20+", "Validated\ndomains")
add_metric_box(slide, Inches(7.3), Inches(1.2), Inches(2), Inches(1), "5M+", "Human survey\nresponses")

# Highlight box
add_highlight_box(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(0.8),
                  "Core insight: Raw LLMs are 25% less accurate than expert forecasters. Calibration against human survey data closes this gap for established topics.",
                  is_dark=True)

# Slide 3: The Problem
slide = add_content_slide(prs, "Raw LLM predictions fail at rates unacceptable for business decisions",
                          "Forecasting Research Institute; Dig Insights validation study (N=500)")

# Two tables side by side
data1 = [
    ["System", "Brier Score", "Gap"],
    ["Superforecasters", "0.081", "—"],
    ["GPT-4.5", "0.101", "+25%"],
    ["GPT-4", "0.131", "+62%"],
    ["Median public", "0.150+", "+85%"]
]
add_table(slide, Inches(0.5), Inches(1.2), Inches(4), data1, [1.8, 1.2, 1])

data2 = [
    ["Prediction Task", "Correlation"],
    ["Known events", "0.85 ✓"],
    ["Future events", "0.50 ⚠"],
    ["New concepts", "0.30 ✗"]
]
add_table(slide, Inches(5), Inches(1.2), Inches(4), data2, [2.5, 1.5])

add_highlight_box(slide, Inches(5), Inches(3.2), Inches(4), Inches(0.7),
                  "The paradox: Synthetic data works for what you already know — and fails at what you need to predict.")

# Slide 4: Accuracy Spectrum
slide = add_content_slide(prs, "Accuracy varies predictably by question type, enabling appropriate use case selection",
                          "Crowdwave accuracy testing (27 test cases, 6 domains)")

# Three zone boxes
add_zone_box(slide, Inches(0.5), Inches(1.2), Inches(2.9), Inches(2),
             "HIGH ACCURACY", "±2-3 pts",
             ["Trust scales", "Awareness (Y/N)", "Party ID", "Bipartisan rankings"],
             "→ Use for decisions", GREEN, LIGHT_GREEN)

add_zone_box(slide, Inches(3.55), Inches(1.2), Inches(2.9), Inches(2),
             "MEDIUM ACCURACY", "±4-5 pts",
             ["Satisfaction (1-5)", "NPS / Recommend", "Concern levels", "Tech comfort"],
             "→ Use for direction", YELLOW, LIGHT_YELLOW)

add_zone_box(slide, Inches(6.6), Inches(1.2), Inches(2.9), Inches(2),
             "LOW ACCURACY", "±8-15 pts",
             ["Purchase intent", "Price sensitivity", "Polarized politics", "Novel behaviors"],
             "→ Validate first", RED, LIGHT_RED)

# Slide 5: High Accuracy Details
slide = add_content_slide(prs, "High-accuracy zone: Stable attitudes with abundant benchmark data",
                          "Gallup (N=13,000+); Pew Research (N=5,000+); AARP 2025 (N=3,838)")

data = [
    ["Question", "Calibrated", "Actual", "Error"],
    ["Trust in scientists", "77%", "77%", "0 pts"],
    ["% Independent", "44%", "45%", "1 pt"],
    ["Smartphone 50+", "89%", "90%", "1 pt"],
    ["Employee engaged", "32%", "31%", "1 pt"]
]
add_table(slide, Inches(0.5), Inches(1.2), Inches(5), data, [2, 1, 1, 1])

# Why box
box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.2), Inches(3.5), Inches(1.8))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_GRAY
box.line.color.rgb = RgbColor(224, 224, 224)

title_txt = slide.shapes.add_textbox(Inches(6.15), Inches(1.3), Inches(3.2), Inches(0.3))
tf = title_txt.text_frame
p = tf.paragraphs[0]
p.text = "Why these work"
set_font(p.runs[0], size=11, bold=True, color=NAVY)

items_txt = slide.shapes.add_textbox(Inches(6.15), Inches(1.6), Inches(3.2), Inches(1.2))
tf = items_txt.text_frame
tf.word_wrap = True
items = ["• Stable attitudes over time", "• Multiple benchmark sources", "• Low emotional volatility", "• Training data aligns with reality"]
for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    set_font(p.runs[0], size=10, color=BLACK)

# Slide 6: NPS Benchmarks
slide = add_content_slide(prs, "Industry-specific calibration required for NPS: actual variance is 30+ points",
                          "Survicate NPS Benchmark 2025 (599 companies, 5.4M responses)")

data = [
    ["Industry", "Median NPS", "B2B", "B2C", "LLM Error"],
    ["Manufacturing", "65", "66", "62", "-25 pts"],
    ["Healthcare", "61", "38", "70", "-20 pts"],
    ["Retail/Ecommerce", "55", "55", "54", "-15 pts"],
    ["Fintech", "46", "—", "—", "-10 pts"],
    ["Software", "30", "29", "47", "+5 pts"]
]
add_table(slide, Inches(0.5), Inches(1.2), Inches(9), data, [2.5, 1.5, 1, 1, 1.5])

add_highlight_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.6),
                  "LLMs assume NPS of 35-40 for all industries. Industry-specific calibration is required for accuracy.")

# Slide 7: Intent Gap & Polarization
slide = add_content_slide(prs, "Low-accuracy zone: Intent requires conversion factors; polarized topics require segmentation",
                          "Meta-analysis (intent gap); Pew Research Feb 2025 (N=5,086)")

# Intent table
data1 = [
    ["Stated Response", "Actual Conversion"],
    ['"Very likely"', "25-35% (×0.30)"],
    ['"Likely"', "10-20% (×0.15)"],
    ['"Might consider"', "3-8% (×0.05)"]
]
add_table(slide, Inches(0.5), Inches(1.2), Inches(4), data1, [2, 2])

# Partisan table
data2 = [
    ["Topic", "Rep", "Dem", "Gap"],
    ["Immigration", "75%", "25%", "50 pts"],
    ["Climate", "25%", "70%", "45 pts"],
    ["Gun violence", "35%", "70%", "35 pts"]
]
add_table(slide, Inches(5), Inches(1.2), Inches(4.5), data2, [1.8, 0.9, 0.9, 0.9])

add_highlight_box(slide, Inches(5), Inches(3.2), Inches(4.5), Inches(0.7),
                  "Never predict a single number for polarized topics. The 'average' represents no one.")

# Slide 8: Bias Patterns
slide = add_content_slide(prs, "Eight documented LLM bias patterns enable systematic correction",
                          "Validated calibrations in CALIBRATION_MEMORY.md")

data = [
    ["Bias Pattern", "Direction", "Correction", "Source"],
    ["Senior tech adoption", "Under-predicts", "×1.30-1.65", "AARP 2025"],
    ["AI concern (general)", "Over-predicts", "×0.90", "Pew/YouGov"],
    ["Status quo preference", "Under-predicts", "+15-20 pts", "Behavioral research"],
    ["Intent-to-action", "Over-predicts", "×0.30-0.55", "Meta-analysis"],
    ["Emotional intensity", "Under-predicts", "×1.20-1.30", "Pet study (N=173)"],
    ["Life satisfaction", "Over-predicts", "-3 to -5 pts", "Gallup 2025"],
    ["Partisan averaging", "Incorrect", "Segment", "Pew 2025"],
    ["Open-end polish", "Over-polished", "20% low-quality", "Industry benchmark"]
]
add_table(slide, Inches(0.5), Inches(1.1), Inches(9), data, [2.2, 1.5, 1.5, 2])

# Slide 9: Validation Results
slide = add_content_slide(prs, "Calibration brings 100% of predictions within 5 points of actual survey results",
                          "Crowdwave validation (27 test cases, 6 domains)")

# Before/after table
data = [
    ["Metric", "Naive LLM", "Calibrated", "Improvement"],
    ["Mean absolute error", "9.1 pts", "1.9 pts", "79%"],
    ["Within 2 pts", "7%", "81%", "+74 pts"],
    ["Within 5 pts", "30%", "100%", "+70 pts"]
]
add_table(slide, Inches(0.5), Inches(1.2), Inches(5), data, [2, 1.2, 1.2, 1.2])

add_highlight_box(slide, Inches(6), Inches(1.2), Inches(3.5), Inches(1),
                  "79% error reduction\nStatistically significant at p < 0.0001", is_dark=True)

# Slide 10: Executive Calibration
slide = add_content_slide(prs, "Executive audiences require role-specific calibration: CHROs are 75% more concerned about AI than CEOs",
                          "Conference Board Global C-Suite Survey 2026 (N=1,732)")

data = [
    ["Concern", "CEO", "CFO", "CHRO", "CMO", "Insight"],
    ["Cyberattacks", "×1.30", "×1.40", "×1.60", "×0.90", "CHROs most concerned"],
    ["AI disruption", "×0.90", "×1.05", "×1.40", "×1.10", "CEOs least concerned"],
    ["Transformation", "×1.50", "×1.15", "×1.70", "×1.40", "CHROs leading change"],
    ["Uncertainty", "×1.35", "×1.50", "×1.50", "×1.25", "CFOs feel it most"]
]
add_table(slide, Inches(0.5), Inches(1.2), Inches(9), data, [1.8, 0.9, 0.9, 0.9, 0.9, 2.2])

add_highlight_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.6),
                  "Generic 'executive' predictions miss role-based variations. Always segment by role when surveying C-suite.")

# Slide 11: Use Cases
slide = add_content_slide(prs, "Use case guidance: Match application to accuracy zone",
                          "Crowdwave accuracy framework")

# Three boxes
add_zone_box(slide, Inches(0.5), Inches(1.1), Inches(2.9), Inches(1.8),
             "HIGH CONFIDENCE", "",
             ["Concept testing", "Audience sizing", "Trend validation", "Priority ranking", "Benchmarking"],
             "", GREEN, LIGHT_GREEN)

add_zone_box(slide, Inches(3.55), Inches(1.1), Inches(2.9), Inches(1.8),
             "MEDIUM CONFIDENCE", "",
             ["Hypothesis generation", "Early screening", "Directional guidance", "", ""],
             "", YELLOW, LIGHT_YELLOW)

add_zone_box(slide, Inches(6.6), Inches(1.1), Inches(2.9), Inches(1.8),
             "VALIDATE FIRST", "",
             ["New products", "Pricing research", "High-stakes decisions", "", ""],
             "", RED, LIGHT_RED)

add_highlight_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.7),
                  "Not recommended without validation: Purchase conversion (use A/B tests), polarized topics (segment by party), novel behaviors (no training data), legal/regulatory evidence")

# Slide 12: Competitive Differentiation
slide = add_content_slide(prs, "Competitive differentiation: Documented accuracy vs. unvalidated claims",
                          "Competitive analysis")

data = [
    ["Capability", "Raw LLM", "Competitors", "Crowdwave"],
    ["Documented accuracy", "None", '"95%" (unvalidated)', "27 test cases ✓"],
    ["Human validation", "None", "Unclear", "5M+ responses ✓"],
    ["Bias corrections", "None", "None documented", "8 patterns ✓"],
    ["Domain calibrations", "None", "Generic", "20+ domains ✓"],
    ["Confidence scoring", "None", "None", "Per-question ✓"],
    ["Known limitations", "None", "None", "Documented ✓"]
]
add_table(slide, Inches(0.5), Inches(1.2), Inches(9), data, [2.2, 1.8, 2.2, 2.2])

add_highlight_box(slide, Inches(0.5), Inches(4), Inches(9), Inches(0.6),
                  "Differentiation: Other vendors claim magic. We document our methodology, show our work, and tell you when NOT to trust the output.")

# Slide 13: Summary
slide = add_content_slide(prs, "Summary: Calibrated predictions deliver 79% error reduction with known accuracy by question type",
                          "Crowdwave Accuracy Framework, February 2026")

# Metrics
add_metric_box(slide, Inches(0.5), Inches(1.1), Inches(2.1), Inches(0.9), "79%", "Error reduction")
add_metric_box(slide, Inches(2.7), Inches(1.1), Inches(2.1), Inches(0.9), "1.9 pts", "Mean error")
add_metric_box(slide, Inches(4.9), Inches(1.1), Inches(2.1), Inches(0.9), "20+", "Domains")
add_metric_box(slide, Inches(7.1), Inches(1.1), Inches(2.1), Inches(0.9), "5M+", "Responses")

# Zones summary
add_zone_box(slide, Inches(0.5), Inches(2.2), Inches(2.9), Inches(1.1),
             "", "±2-3 pts", ["Trust, awareness"], "→ Decisions", GREEN, LIGHT_GREEN)
add_zone_box(slide, Inches(3.55), Inches(2.2), Inches(2.9), Inches(1.1),
             "", "±4-5 pts", ["Satisfaction, NPS"], "→ Direction", YELLOW, LIGHT_YELLOW)
add_zone_box(slide, Inches(6.6), Inches(2.2), Inches(2.9), Inches(1.1),
             "", "±8-15 pts", ["Intent, polarized"], "→ Validate", RED, LIGHT_RED)

add_highlight_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.6),
                  "Documented accuracy. Known limits. Transparent methodology.", is_dark=True)

# Slide 14: End slide
add_title_slide(prs, "Crowdwave", "Documented accuracy. Known limits. Transparent methodology.")

# Save
prs.save('CROWDWAVE_PROFESSIONAL.pptx')
print("Created CROWDWAVE_PROFESSIONAL.pptx")
