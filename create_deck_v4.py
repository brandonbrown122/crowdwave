"""
CrowdWave Consulting Deck V4 - With Charts
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

# Colors
NAVY = RGBColor(0, 51, 102)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 100, 100)
LIGHT_GRAY = RGBColor(240, 240, 240)
GREEN = RGBColor(34, 139, 34)
RED = RGBColor(178, 34, 34)
BLUE = RGBColor(70, 130, 180)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Navy background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_func):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "CrowdWave | Confidential | February 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    
    content_func(slide)
    return slide

def add_bar_chart(slide, x, y, w, h, categories, series_data, title=""):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data.items():
        chart_data.add_series(name, values)
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    ).chart
    
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    
    return chart

def add_horizontal_bar(slide, x, y, w, h, categories, values, colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Values', values)
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    ).chart
    
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(11)
    data_labels.font.bold = True
    
    return chart

def add_text_box(slide, x, y, w, h, text, size=14, bold=False, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return box

# === SLIDES ===

# Slide 1: Title
add_title_slide("CrowdWave", "Accurate consumer insights in minutes, not months\n\nFebruary 2026")

# Slide 2: Executive Summary
def exec_summary(slide):
    # SCR boxes
    labels = ["SITUATION", "COMPLICATION", "RESOLUTION"]
    texts = [
        "Market research takes 4-6 weeks and $25K+ per study. By the time you have answers, the market has moved.",
        "Competitors running 20 concept tests per quarter will outlearn teams running 2-3. Speed is the new competitive moat.",
        "CrowdWave delivers 95% directional accuracy in minutes. Validated against Pew, Gallup, AARP — 2-point average error. Test 10x more, kill losers instantly, validate only winners."
    ]
    colors = [BLUE, RGBColor(180, 100, 50), GREEN]
    
    for i, (label, text, color) in enumerate(zip(labels, texts, colors)):
        y = 1.5 + i * 1.8
        # Label box
        label_box = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(2), Inches(0.5))
        label_box.fill.solid()
        label_box.fill.fore_color.rgb = color
        label_box.line.fill.background()
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Text
        add_text_box(slide, 2.7, y, 10, 1.5, text, size=16)

add_content_slide("Executive Summary", exec_summary)

# Slide 3: The Speed Gap - Bar Chart
def speed_gap(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "At $25K per study, your budget buys 3 tests. Competitors using simulation run 50.", size=16, bold=True, color=NAVY)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Traditional\nResearch', 'CrowdWave\nSimulation']
    chart_data.add_series('Concept Tests per Quarter', (3, 50))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2.2), Inches(5.5), Inches(4), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(18)
    plot.data_labels.font.bold = True
    
    # Side stats
    stats = [
        ("$25,000+", "Cost per traditional study"),
        ("4-6 weeks", "Time to insight"),
        ("~$0", "Marginal cost with simulation"),
        ("Minutes", "Time to first results")
    ]
    for i, (num, label) in enumerate(stats):
        y = 2.4 + i * 1.1
        add_text_box(slide, 7.5, y, 2, 0.5, num, size=28, bold=True, color=NAVY)
        add_text_box(slide, 9.7, y + 0.1, 3, 0.5, label, size=13)

add_content_slide("Your research budget buys 3 studies — competitors are testing 50", speed_gap)

# Slide 4: Accuracy Proof - Bar Chart
def accuracy_proof(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "27 blind tests against Pew, Gallup, AARP. Mean error: 1.9 points.", size=16, bold=True, color=NAVY)
    
    # Prediction vs Actual bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Smartphone\n50+', 'Independent\nvoters', 'Trust in\nscientists', 'AI\nconcern', 'Mfg\nNPS']
    chart_data.add_series('CrowdWave Prediction', (89, 44, 77, 50, 64))
    chart_data.add_series('Actual Result', (90, 45, 77, 48, 65))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(8), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # Error callouts
    add_text_box(slide, 9, 2.2, 3.8, 3, 
        "Error by test:\n\n• Smartphone 50+: 1 pt\n• Independent voters: 1 pt\n• Trust scientists: 0 pts\n• AI concern: 2 pts\n• Mfg NPS: 1 pt\n\n✓ 100% within 5 points", 
        size=14)

add_content_slide("We predicted real consumer behavior within 2 points — blindly", accuracy_proof)

# Slide 5: Calibration Impact - Before/After
def calibration_impact(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "Raw AI averages 9-point error. Calibration cuts it to 1.9 points.", size=16, bold=True, color=NAVY)
    
    # Before/After comparison
    chart_data = CategoryChartData()
    chart_data.categories = ['Mean Absolute\nError (pts)', 'Within 2 pts\n(%)', 'Within 5 pts\n(%)']
    chart_data.add_series('Raw LLM', (9.1, 7, 30))
    chart_data.add_series('Calibrated CrowdWave', (1.9, 81, 100))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(7), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # Key insight box
    box = slide.shapes.add_shape(1, Inches(8), Inches(2.5), Inches(4.8), Inches(3))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    
    add_text_box(slide, 8.2, 2.7, 4.4, 0.5, "79%", size=48, bold=True, color=NAVY)
    add_text_box(slide, 8.2, 3.5, 4.4, 0.4, "Error reduction", size=18, color=NAVY)
    add_text_box(slide, 8.2, 4.1, 4.4, 1.2, 
        "Built from:\n• 8 documented bias patterns\n• 20+ domain calibrations\n• 5M+ human survey responses", size=13)

add_content_slide("Raw AI predictions fail — calibration makes them reliable", calibration_impact)

# Slide 6: Accuracy Zones - Visual Matrix
def accuracy_zones(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "Match the tool to the question type for predictable accuracy.", size=16, bold=True, color=NAVY)
    
    # Three zones with horizontal bars showing error ranges
    zones = [
        ("HIGH ACCURACY", "±2-3 pts", "Trust, awareness, party ID, demographics", GREEN, ["Trust scales", "Awareness", "Demographics"]),
        ("MEDIUM ACCURACY", "±4-5 pts", "Satisfaction, NPS, concern levels", RGBColor(200, 150, 50), ["Satisfaction", "NPS", "Concerns"]),
        ("LOW ACCURACY", "±8-15 pts", "Purchase intent, price sensitivity, polarized", RED, ["Intent", "Pricing", "Polarized"])
    ]
    
    for i, (label, error, desc, color, items) in enumerate(zones):
        y = 2 + i * 1.7
        
        # Zone label
        label_box = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(2.5), Inches(0.6))
        label_box.fill.solid()
        label_box.fill.fore_color.rgb = color
        label_box.line.fill.background()
        tf = label_box.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Error
        add_text_box(slide, 3.2, y, 1.5, 0.6, error, size=16, bold=True, color=color)
        
        # Description
        add_text_box(slide, 4.8, y + 0.1, 4, 0.5, desc, size=13)
        
        # Recommendation icons
        rec = "✓ Use for decisions" if i == 0 else ("✓ Use for direction" if i == 1 else "⚠ Validate first")
        add_text_box(slide, 9.5, y + 0.1, 3.5, 0.5, rec, size=13, bold=True)
    
    # Example callouts
    add_text_box(slide, 0.5, 5.8, 12, 1, 
        "Examples:  'Which 3 of 10 concepts resonate?' = Simulation alone  |  'How much would they pay?' = Validate  |  'Immigration views?' = Segment by party",
        size=12, color=GRAY)

add_content_slide("Accuracy is predictable by question type", accuracy_zones)

# Slide 7: Senior Tech Bias - Correction Chart
def senior_bias(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "LLMs underestimate adults 60+ by 25% on technology adoption.", size=16, bold=True, color=NAVY)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Smartphone\nownership\n(50+)', 'Daily internet\nuse (60+)', 'Video\nstreaming\n(70+)']
    chart_data.add_series('Raw LLM', (72, 60, 40))
    chart_data.add_series('Calibrated', (89, 82, 65))
    chart_data.add_series('Actual', (90, 83, 64))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(7), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # Correction factors
    add_text_box(slide, 8, 2.2, 4.8, 0.4, "Correction factors:", size=14, bold=True, color=NAVY)
    factors = [
        ("50-69:", "×1.30"),
        ("70-79:", "×1.40"),
        ("80+:", "×1.50")
    ]
    for i, (age, factor) in enumerate(factors):
        add_text_box(slide, 8, 2.8 + i * 0.6, 2, 0.5, age, size=14)
        add_text_box(slide, 9.5, 2.8 + i * 0.6, 1.5, 0.5, factor, size=14, bold=True, color=GREEN)
    
    add_text_box(slide, 8, 4.8, 4.8, 1.5, 
        "Why it happens:\nLLM training over-represents stereotypes. Reality has shifted — 90% of 50+ own smartphones.\n\nSource: AARP Tech Trends 2025 (N=3,838)", 
        size=11, color=GRAY)

add_content_slide("LLMs systematically underestimate seniors — we found the fix", senior_bias)

# Slide 8: Polarization Chart
def polarization(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "The 'average American' is a fiction on partisan topics.", size=16, bold=True, color=NAVY)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Immigration\nconcern', 'Climate\nconcern', 'Gun violence\nconcern']
    chart_data.add_series('Republican', (75, 25, 35))
    chart_data.add_series('Democrat', (25, 70, 70))
    chart_data.add_series('"Average"', (48, 45, 52))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(7.5), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # Gap callouts
    add_text_box(slide, 8.3, 2.2, 4.5, 0.4, "Partisan gaps:", size=14, bold=True, color=NAVY)
    gaps = [("Immigration:", "50 pts"), ("Climate:", "45 pts"), ("Gun violence:", "35 pts")]
    for i, (issue, gap) in enumerate(gaps):
        add_text_box(slide, 8.3, 2.8 + i * 0.6, 2, 0.5, issue, size=14)
        add_text_box(slide, 10.3, 2.8 + i * 0.6, 1.5, 0.5, gap, size=14, bold=True, color=RED)
    
    # Rule box
    box = slide.shapes.add_shape(1, Inches(8.3), Inches(4.6), Inches(4.5), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    add_text_box(slide, 8.5, 4.8, 4.2, 1.2, 
        "Rule: Never report a single number on polarized topics.\n\nCrowdWave enforces segmentation automatically.", size=12, bold=True)

add_content_slide("Political topics require segmentation — or you'll miss by 50 points", polarization)

# Slide 9: Intent Gap
def intent_gap(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "Stated intent overstates actual behavior by 3-5x.", size=16, bold=True, color=NAVY)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['"Definitely\nwill buy"', '"Probably\nwill buy"', '"Might\nconsider"']
    chart_data.add_series('Stated Intent', (85, 55, 35))
    chart_data.add_series('Actual Conversion', (30, 15, 5))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(6.5), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # Correction factors
    add_text_box(slide, 7.5, 2.2, 5.3, 0.4, "Conversion factors:", size=14, bold=True, color=NAVY)
    factors = [
        ('"Definitely"', '→', '×0.30'),
        ('"Probably"', '→', '×0.15'),
        ('"Might"', '→', '×0.05')
    ]
    for i, (stated, arrow, factor) in enumerate(factors):
        add_text_box(slide, 7.5, 2.8 + i * 0.7, 2, 0.5, stated, size=14)
        add_text_box(slide, 9.5, 2.8 + i * 0.7, 0.5, 0.5, arrow, size=14)
        add_text_box(slide, 10, 2.8 + i * 0.7, 1.5, 0.5, factor, size=16, bold=True, color=NAVY)
    
    add_text_box(slide, 7.5, 5, 5.3, 1.5, 
        "CrowdWave applies these corrections automatically when purchase intent questions are detected.\n\nFor pricing: always validate with behavioral data.", 
        size=12, color=GRAY)

add_content_slide("Purchase intent overstates reality by 3-5x — we apply corrections", intent_gap)

# Slide 10: C-Suite Calibration
def csuite(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "CHROs are 75% more worried about AI disruption than CEOs.", size=16, bold=True, color=NAVY)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Cyberattacks', 'AI disruption', 'Economic\nuncertainty']
    chart_data.add_series('CEO', (130, 90, 135))
    chart_data.add_series('CFO', (140, 105, 150))
    chart_data.add_series('CHRO', (160, 140, 150))
    chart_data.add_series('CMO', (90, 110, 125))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(7.5), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    add_text_box(slide, 0.5, 6.3, 7.5, 0.5, "Index: 100 = generic 'executive' baseline", size=10, color=GRAY)
    
    # Insight box
    box = slide.shapes.add_shape(1, Inches(8.3), Inches(2.2), Inches(4.5), Inches(3))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    add_text_box(slide, 8.5, 2.4, 4.2, 2.6, 
        "Key insight:\n\nGeneric 'executive' predictions miss role variation by 40+ points.\n\n• CHROs: +40% AI concern\n• CMOs: -10% cyber concern\n• CEOs: +35% economy focus\n\nSpecify the role. Generic wastes accuracy.", 
        size=12)

add_content_slide("C-suite predictions require role-specific calibration", csuite)

# Slide 11: Industry NPS
def industry_nps(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "Industry NPS ranges from 30 to 65 — LLMs assume everyone is at 40.", size=16, bold=True, color=NAVY)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Manufacturing', 'Healthcare', 'Retail', 'Fintech', 'Software']
    chart_data.add_series('Actual NPS', (65, 61, 55, 46, 30))
    chart_data.add_series('LLM Prediction', (40, 40, 40, 40, 40))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(7.5), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # Error callouts
    add_text_box(slide, 8.3, 2.2, 4.5, 0.4, "LLM error by industry:", size=14, bold=True, color=NAVY)
    errors = [
        ("Manufacturing:", "-25 pts"),
        ("Healthcare:", "-21 pts"),
        ("Retail:", "-15 pts"),
        ("Fintech:", "-6 pts"),
        ("Software:", "+10 pts")
    ]
    for i, (ind, err) in enumerate(errors):
        add_text_box(slide, 8.3, 2.8 + i * 0.55, 2.2, 0.5, ind, size=13)
        color = RED if err.startswith("-") else RGBColor(200, 150, 50)
        add_text_box(slide, 10.5, 2.8 + i * 0.55, 1.5, 0.5, err, size=13, bold=True, color=color)
    
    add_text_box(slide, 8.3, 5.6, 4.5, 0.8, 
        "Our fix: Industry-specific baselines. Manufacturing starts at 65.\n\nSource: Survicate 2025 (N=5.4M)", 
        size=11, color=GRAY)

add_content_slide("Industry NPS varies 35 points — raw AI misses this completely", industry_nps)

# Slide 12: Economics Impact
def economics(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "10x more learning at 1/100th the cost. Compounding advantage.", size=16, bold=True, color=NAVY)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Time to\ninsight', 'Cost per\nconcept', 'Tests per\nquarter', 'Iterations']
    chart_data.add_series('Traditional (indexed)', (100, 100, 100, 100))
    chart_data.add_series('CrowdWave (indexed)', (0.1, 0.1, 2000, 5000))
    
    # Use log scale conceptually - show multipliers instead
    chart_data2 = CategoryChartData()
    chart_data2.categories = ['Time to insight', 'Cost per concept', 'Tests per quarter', 'Iterations possible']
    chart_data2.add_series('Improvement Multiple', (1000, 100, 20, 50))
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(2.2), Inches(6.5), Inches(4), chart_data2
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(14)
    plot.data_labels.font.bold = True
    
    # Side narrative
    add_text_box(slide, 7.5, 2.2, 5.3, 0.4, "The compounding effect:", size=14, bold=True, color=NAVY)
    steps = [
        "Week 1: Simulate 20 concepts, kill 15",
        "Week 2: Iterate on 5 survivors",
        "Week 3: Validate top 2 ($50K)",
        "Week 4: Launch with confidence"
    ]
    for i, step in enumerate(steps):
        add_text_box(slide, 7.5, 2.8 + i * 0.6, 5.3, 0.5, step, size=13)
    
    add_text_box(slide, 7.5, 5.3, 5.3, 1, 
        "Traditional: Test 2 concepts in 6 weeks. Hope you picked right.", 
        size=13, bold=True, color=GRAY)

add_content_slide("Simulation transforms research economics", economics)

# Slide 13: Decision Framework - 2x2 Matrix
def decision_matrix(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "Match simulation confidence to decision stakes.", size=16, bold=True, color=NAVY)
    
    # Draw 2x2 matrix
    # Headers
    add_text_box(slide, 4.5, 1.9, 3.5, 0.5, "LOW STAKES", size=14, bold=True, color=NAVY)
    add_text_box(slide, 8.5, 1.9, 3.5, 0.5, "HIGH STAKES", size=14, bold=True, color=NAVY)
    
    # Row labels
    add_text_box(slide, 0.7, 2.8, 2.5, 0.8, "HIGH\nACCURACY", size=12, bold=True, color=NAVY)
    add_text_box(slide, 0.7, 4.3, 2.5, 0.8, "MEDIUM\nACCURACY", size=12, bold=True, color=NAVY)
    add_text_box(slide, 0.7, 5.8, 2.5, 0.8, "LOW\nACCURACY", size=12, bold=True, color=NAVY)
    
    # Matrix cells
    cells = [
        (3.5, 2.5, GREEN, "✓ Simulation only"),
        (7.5, 2.5, GREEN, "✓ Simulation +\nspot validation"),
        (3.5, 4, RGBColor(200, 180, 50), "✓ Directional use"),
        (7.5, 4, RGBColor(200, 150, 50), "⚠ Validate before\nmajor spend"),
        (3.5, 5.5, RGBColor(200, 150, 50), "⚠ Directional only"),
        (7.5, 5.5, RED, "✗ Always validate")
    ]
    
    for x, y, color, text in cells:
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.5), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = color
        box.line.width = Pt(3)
        add_text_box(slide, x + 0.2, y + 0.3, 3.2, 1, text, size=13, bold=True)
    
    # Thresholds
    add_text_box(slide, 0.5, 6.6, 12, 0.5, 
        "Thresholds:  <$100K → Simulation sufficient  |  $100K-$1M → Validate finalists  |  >$1M → Simulation screens, humans decide",
        size=12, color=GRAY)

add_content_slide("Decision framework: Match confidence to stakes", decision_matrix)

# Slide 14: What We Built
def what_we_built(slide):
    add_text_box(slide, 0.5, 1.4, 12, 0.5, "20+ validated domains, 8 bias corrections, 5M+ human responses.", size=16, bold=True, color=NAVY)
    
    # Component bars
    components = [
        ("Validated domains", 20),
        ("Bias corrections", 8),
        ("Calibration factors", 100),
        ("Test cases", 27)
    ]
    
    chart_data = CategoryChartData()
    chart_data.categories = [c[0] for c in components]
    chart_data.add_series('Count', [c[1] for c in components])
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(2), Inches(6), Inches(3.5), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(14)
    plot.data_labels.font.bold = True
    
    # Human data callout
    box = slide.shapes.add_shape(1, Inches(7), Inches(2), Inches(5.8), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    add_text_box(slide, 7.3, 2.2, 2, 0.8, "5M+", size=36, bold=True, color=WHITE)
    add_text_box(slide, 9.5, 2.3, 3, 0.8, "Human survey\nresponses", size=16, color=WHITE)
    
    # Source tiers
    add_text_box(slide, 7, 3.5, 5.8, 0.4, "Source quality tiers:", size=13, bold=True, color=NAVY)
    tiers = [
        "Tier 1: Pew, Gallup, AARP (probability samples, N>1K)",
        "Tier 2: McKinsey, Deloitte, Conference Board",
        "Tier 3: YouGov, Harris Poll (directional)"
    ]
    for i, tier in enumerate(tiers):
        add_text_box(slide, 7, 4 + i * 0.5, 5.8, 0.5, tier, size=11)
    
    # Domains list
    add_text_box(slide, 0.5, 5.8, 12, 0.8, 
        "Domains: Trust, technology adoption, NPS by industry, executive attitudes, consumer concerns, travel/hospitality, healthcare, political identity", 
        size=11, color=GRAY)

add_content_slide("System foundation: Validated calibrations at scale", what_we_built)

# Slide 15: Three Actions
def three_actions(slide):
    actions = [
        ("1", "Integrate simulation into every research project", 
         "Simulate first. Screen concepts, kill losers. Then decide what needs validation."),
        ("2", "Set decision thresholds by stakes",
         "Screening → Simulation only. Major campaigns → Validate finalists. Pricing → Always validate."),
        ("3", "Track and compound accuracy",
         "Log predictions vs. outcomes. Feed misses back. Calibration improves continuously.")
    ]
    
    for i, (num, title, desc) in enumerate(actions):
        y = 1.6 + i * 1.8
        
        # Number circle
        circle = slide.shapes.add_shape(9, Inches(0.5), Inches(y), Inches(0.7), Inches(0.7))  # Oval
        circle.fill.solid()
        circle.fill.fore_color.rgb = NAVY
        circle.line.fill.background()
        add_text_box(slide, 0.6, y + 0.1, 0.5, 0.5, num, size=24, bold=True, color=WHITE)
        
        # Title and description
        add_text_box(slide, 1.5, y, 11, 0.5, title, size=18, bold=True, color=NAVY)
        add_text_box(slide, 1.5, y + 0.6, 11, 0.8, desc, size=14)
    
    # Closing statement
    box = slide.shapes.add_shape(1, Inches(0.5), Inches(5.8), Inches(12.333), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    add_text_box(slide, 0.7, 6, 12, 0.7, 
        "The question isn't whether to use simulation — it's how much ground you'll lose to competitors who start first.",
        size=16, bold=True, color=NAVY)

add_content_slide("Three actions to capture the speed advantage", three_actions)

# Slide 16: Closing
add_title_slide("CrowdWave", "Documented accuracy. Known limits. Transparent methodology.\n\nFebruary 2026")

# Save
prs.save('C:/Users/brand/clawd-crowdwave/CROWDWAVE_DECK_V4.pptx')
print("Created CROWDWAVE_DECK_V4.pptx")
